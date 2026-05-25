#!/usr/bin/env python3
"""Build Fortiva Technical Overview PowerPoint from repo-accurate content."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent
IMAGES = ROOT / "images"
OUTPUT = ROOT / "Fortiva-Technical-Overview.pptx"

NAVY = RGBColor(0x0B, 0x1A, 0x2E)
TEAL = RGBColor(0x00, 0xB4, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)


def set_slide_bg(slide, rgb=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    img = IMAGES / "fortiva-title-cover.png"
    if img.exists():
        slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEAL
        p2.space_before = Pt(12)


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, TEAL)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, image_name=None, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)

    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = TEAL

    content_left = Inches(0.5)
    content_top = Inches(1.1)
    content_width = Inches(9.0)
    if image_name:
        img_path = IMAGES / image_name
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(5.0), Inches(1.1), width=Inches(4.5))
            content_width = Inches(4.3)

    body = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            level, text = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(14 if (isinstance(item, tuple) and item[0] > 0) else 15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_table_slide(prs, title, headers, rows, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = TEAL

    cols, row_count = len(headers), len(rows) + 1
    table = slide.shapes.add_table(row_count, cols, Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.4 * row_count)).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = TEAL
            p.font.size = Pt(12)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.size = Pt(11)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_code_slide(prs, title, code_lines, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = TEAL

    code_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9.0), Inches(5.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x22)
    code_box.line.color.rgb = TEAL

    tf = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.1)).text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Fortiva Technical Overview",
        "Zero-Knowledge · Local-First · Windows-Native  |  icmclab studio",
    )

    add_content_slide(prs, "What Fortiva Is", [
        "Fortiva is a Windows-native password manager built on WinUI 3 and .NET 8.",
        "Zero-knowledge: your master password is never stored. Only derived keys and encrypted blobs exist on disk.",
        "Local-first: no cloud sync, no telemetry, no background service when the app is closed.",
        "Living off the land: AES-256-GCM via Windows CNG, DPAPI for local metadata, Argon2id for KDF.",
        "Three shipping applications plus CLI tools and a browser extension bridge.",
    ], notes="Audience may include IT, security reviewers, and developers. Emphasize offline operation.")

    add_content_slide(prs, "Solution Components", [
        "Fortiva.Core — shared library: vault format, crypto, policy, audit, browser bridge server.",
        "Fortiva.Personal — free WinUI app for individuals (%APPDATA% vault).",
        "Fortiva.Enterprise — licensed WinUI app with IT policy, seats, shared vaults.",
        "Fortiva.Admin — IT console (requires Administrator) for license, policy, shared vault config.",
        "Fortiva.LicenseTool — CLI to generate RSA keys and sign enterprise licenses.",
        "Fortiva.BrowserBridge.Host — native messaging host for Chromium/Edge extension.",
        "extension/ — browser extension (local-only credential fill via named pipes).",
    ])

    add_content_slide(prs, "Three Editions (Who Uses What)", [
        "Personal: individual users, vault at %APPDATA%\\Fortiva\\vault.fva, optional auto-update check.",
        "Enterprise: business users, vault default %PROGRAMDATA%\\Fortiva, license + policy required.",
        "Admin: IT staff only, does not store end-user passwords; manages license.dat and policies.json.",
        "Enterprise users pick shared vault paths configured by Admin in Settings.",
        "All editions share the same .fva vault cryptography and Fortiva.Core engine.",
    ], image_name="fortiva-editions.png")

    add_section_slide(prs, "Part 1 — Vault Cryptography")

    add_content_slide(prs, "Zero-Knowledge Model (Explicit)", [
        "YOU know: master password (never written to disk by Fortiva).",
        "DISK contains: encrypted vault file, wrapped vault key, salts, KDF parameters, DPAPI metadata.",
        "MEMORY (while unlocked): Master Key (MK) and Vault Key (VK) — wiped on lock/panic.",
        "Fortiva cannot reset your master password. Loss = no recovery unless you have a backup export.",
        "IT admins can set policy and license but cannot decrypt vault entries without the master password.",
    ], image_name="fortiva-key-hierarchy.png", notes="Walk through that MK is derived on every unlock and zeroed after lock.")

    add_table_slide(prs, "Argon2id Parameters (From Source Code)", [
        "Profile", "Memory (KB)", "Iterations", "Parallelism", "Typical use"
    ], [
        ("PersonalDefault", "65,536 (64 MB)", "3", "4", "Fortiva Personal new vaults"),
        ("EnterpriseMinimum (policy floor)", "131,072 (128 MB)", "4", "4", "Strict enterprise policy default"),
        ("Paranoia", "262,144 (256 MB)", "5", "4", "Paranoia security level / mandatory policy"),
    ], notes="Implementation: Isopoh.Cryptography.Argon2, Argon2Type.HybridAddressing, Argon2Version.Nineteen.")

    add_content_slide(prs, "Key Hierarchy (Step by Step)", [
        "1. User enters master password at unlock.",
        "2. Argon2id derives 32-byte Master Key (MK) using per-vault salt + stored KDF parameters.",
        "3. MK unwraps Vault Key (VK) from header field wrapped_vault_key (AES-256-GCM).",
        "4. VK decrypts entries blob and integrity log blob inside vault.fva.",
        "5. On lock: MK/VK zeroed via CryptographicOperations.ZeroMemory; bridge token cleared.",
        "Security levels: Standard, Enhanced, Paranoia — Paranoia uses stronger Argon2 profile.",
    ], image_name="fortiva-key-hierarchy.png")

    add_content_slide(prs, "Symmetric Encryption — Windows CNG", [
        "Algorithm: AES-256-GCM (authenticated encryption) through Windows CNG.",
        "Header MAC: canonical header bytes authenticated under VK.",
        "Payload: entries JSON + integrity log encrypted as separate GCM blobs.",
        "Export backup: separate Argon2id + AES-GCM wrapper with user-chosen export password.",
        "No custom crypto primitives — platform CNG and vetted Argon2 library (Konscious/Isopoh).",
    ])

    add_content_slide(prs, "Vault File Layout (.fva)", [
        "Magic string FORTIVA then versioned header fields.",
        "Header stores: vault_id (GUID), created/modified UTC, revision_counter, security_level_counter.",
        "KDF blob, salt, wrapped_vault_key, header_mac.",
        "Encrypted entries: titles, usernames, passwords, URLs, notes, TOTP secrets (Enterprise).",
        "Encrypted integrity log: add/update/delete/import events for consistency checks.",
        "Max vault file size enforced (VaultConstants.MaxVaultFileBytes).",
    ], image_name="fortiva-vault-format.png")

    add_content_slide(prs, "Atomic Save + Snapshot Rotation", [
        "Write protocol: write vault.fva.tmp → flush to disk → File.Replace into vault.fva.",
        "After each save, rotate snapshots: snapshot5 deleted, 4→5, 3→4, 2→3, 1→2, current→snapshot1.",
        "Default SnapshotCount = 5 (vault.fva.snapshot1 … snapshot5).",
        "Restore: unlock snapshot with password, verify, optionally replace live vault.",
        "Protects against partial writes and some corruption scenarios.",
    ])

    add_content_slide(prs, "Rollback & Paranoia Protection", [
        "Separate file local.state — DPAPI-protected metadata (Personal: CurrentUser, Enterprise: LocalMachine scope for vault dir).",
        "Tracks: max security_level seen, last vault_id, last modified time, revision_counter.",
        "If vault revision_counter decreases or security_level downgrades suspiciously → rollback detected.",
        "Default behavior: vault opens READ-ONLY until user explicitly confirms rollback.",
        "Paranoia Mode (user or policy): stricter posture; enterprise can mandate via MandatoryParanoiaMode.",
    ])

    add_section_slide(prs, "Part 2 — Runtime & Data Locations")

    add_table_slide(prs, "On-Disk Locations (Personal vs Enterprise)", [
        "Asset", "Personal", "Enterprise"
    ], [
        ("Vault file", "%APPDATA%\\Fortiva\\vault.fva", "%PROGRAMDATA%\\Fortiva\\vault.fva (default)"),
        ("local.state (rollback)", "Same folder as vault", "Same folder as vault"),
        ("Windows Hello blobs", "%APPDATA%\\Fortiva\\hello.*", "%LOCALAPPDATA%\\FortivaEnterprise\\Hello\\"),
        ("Audit logs", "%LOCALAPPDATA%\\FortivaPersonal\\audit\\", "%PROGRAMDATA%\\Fortiva\\audit\\"),
        ("License / policy", "N/A", "%PROGRAMDATA%\\Fortiva\\license.dat, policies.json"),
        ("Seat registry", "N/A", "%PROGRAMDATA%\\Fortiva\\seats.dat"),
        ("Shared vault config", "N/A", "%PROGRAMDATA%\\Fortiva\\shared-vaults.json"),
        ("Portable USB vault", "User-selected path\\Fortiva\\vault.fva", "Blocked by typical enterprise policy"),
    ])

    add_content_slide(prs, "Unlock Flow (End User)", [
        "App starts → checks vault.fva exists → Unlock page or Onboarding.",
        "Password path: Argon2id → unwrap VK → decrypt payload → StartInfrastructure().",
        "Hello path: UserConsentVerifier biometric/PIN gate → load DPAPI hello.keyprotect → MK.",
        "Enterprise: valid license + seat available checked before unlock succeeds.",
        "While unlocked: auto-lock timer, browser bridge pipes, clipboard policy enforced.",
        "Lock/Panic: dispose session keys, stop bridge, clear in-memory bridge token.",
    ])

    add_content_slide(prs, "Windows Hello (How It Works — No Assumptions)", [
        "Hello does NOT replace the master password cryptographically for vault creation.",
        "Setup flow: user verifies master password → app copies MK from session → stores hello bundle.",
        "Bundle format v3 (FTWH): random unlock key + MK wrapped with AES-GCM, then DPAPI-protected.",
        "DPAPI scope: always CurrentUser (never LocalMachine) to prevent cross-user decrypt on shared PCs.",
        "Unlock: UserConsentVerifier must succeed, then TryLoadMasterKey(helloVerified: true).",
        "Limitation (documented): same-user malware could extract MK from DPAPI blob — Hello is a consent gate.",
    ])

    add_content_slide(prs, "Browser Bridge Architecture", [
        "Extension never reads vault file directly.",
        "Native host Fortiva.BrowserBridge.Host.exe receives JSON via Chrome native messaging.",
        "Host requests session token from secured pipe Fortiva.Bridge.Token (in-memory, same user ACL).",
        "Host connects to named pipe Fortiva.BrowserBridge with token + domain payload.",
        "Fortiva app returns username/password for matching entry only while vault unlocked.",
        "Client validation: bridge executable path checked against install roots (fail closed in release).",
    ], image_name="fortiva-browser-bridge.png")

    add_section_slide(prs, "Part 3 — Enterprise Licensing & Policy")

    add_content_slide(prs, "License Document Fields", [
        "Edition (default Enterprise), CompanyName, ExpiresAt (UTC).",
        "FeatureFlags: vault, policy, audit, shared_vaults.",
        "MaxSeats: enforced via LicenseSeatRegistry on enterprise unlock.",
        "Signature: RSA-PKCS1-SHA256 over canonical JSON payload.",
        "Storage: license.dat = DPAPI-protected (LocalMachine) under %PROGRAMDATA%\\Fortiva\\.",
        "Portable distribution: signed .json file importable via Admin or TryImportFromFile.",
    ], image_name="fortiva-licensing-workflow.png")

    add_code_slide(prs, "LicenseTool — Step 1: Build the CLI", [
        "cd C:\\Repo\\Github\\Fortiva",
        "dotnet build src\\Fortiva.LicenseTool\\Fortiva.LicenseTool.csproj -c Release",
        "",
        "# Run from repo (development):",
        "dotnet run --project src\\Fortiva.LicenseTool -- generate-key",
    ], notes="Release deployments ship LicenseTool beside Admin or as separate CLI.")

    add_code_slide(prs, "LicenseTool — Step 2: Generate RSA Key Pair (First Time Only)", [
        "dotnet run --project src\\Fortiva.LicenseTool -- generate-key",
        "",
        "# Output:",
        "#  PUBLIC KEY  → embed in LicenseVerifier.EmbeddedPublicKeyXml before production build",
        "#  PRIVATE KEY → save to secure file, e.g. C:\\Secure\\fortiva-private.xml",
        "",
        "# NEVER commit private key to git.",
        "# Release Enterprise/Admin builds fail if dev public key still embedded (unless DEBUG).",
    ], notes="Production: replace EmbeddedPublicKeyXml in LicenseDocument.cs / LicenseVerifier with your org public key.")

    add_code_slide(prs, "LicenseTool — Step 3: Sign a License", [
        "# Syntax (from Program.cs):",
        "# sign <company-name> <days-valid> <private-key.xml>",
        "",
        "dotnet run --project src\\Fortiva.LicenseTool -- sign \"Acme Corp\" 365 C:\\Secure\\fortiva-private.xml",
        "",
        "# Writes:",
        "#   %PROGRAMDATA%\\Fortiva\\license.dat  (DPAPI-protected, machine scope)",
        "#   .\\fortiva-license-acme-corp.json     (portable JSON for Admin import)",
        "",
        "# Default MaxSeats in signer: 100 (edit LicenseTool Program.cs to customize).",
    ])

    add_code_slide(prs, "LicenseTool — Step 4: Verify a License", [
        "# Verify portable JSON:",
        "dotnet run --project src\\Fortiva.LicenseTool -- verify .\\fortiva-license-acme-corp.json",
        "",
        "# Verify installed license.dat (no args — loads from %PROGRAMDATA%):",
        "dotnet run --project src\\Fortiva.LicenseTool -- verify",
        "",
        "# Output shows: Company, Edition, Expires, Signature VALID/INVALID, Active status.",
    ])

    add_content_slide(prs, "Deploying License to Clients", [
        "Option A — Admin Console: run Fortiva Admin as Administrator → License tab → Import License → pick .json or .dat.",
        "Option B — MDM/script: copy signed license.dat to %PROGRAMDATA%\\Fortiva\\ (requires admin).",
        "Option C — Intune: wrap installer + config (see packaging/intune/README.md).",
        "Enterprise client loads license on startup; invalid/expired license blocks vault create/unlock.",
        "Admin shows seat usage: CountActiveSeats() / MaxSeats from seats.dat.",
    ], image_name="fortiva-licensing-workflow.png")

    add_content_slide(prs, "Enterprise Policy Engine", [
        "policies.json — DPAPI LocalMachine protected at %PROGRAMDATA%\\Fortiva\\.",
        "Fields: Min Argon2 memory/iters/parallelism, MaxAutoLockSeconds, clipboard mode/clear, ExportMode.",
        "PortableModeAllowed, MandatoryParanoiaMode, MandatoryWindowsHello, TotpEnabled.",
        "HKLM registry overrides: SOFTWARE\\Fortiva\\Enterprise\\Policy (documented for Intune).",
        "PolicyEnforcer applies floors in UI and Core (KDF, export, clipboard, security level on save).",
        "Users cannot weaken below IT baseline; violations logged to audit trail.",
    ])

    add_content_slide(prs, "Shared Vaults & Seat Enforcement", [
        "Admin → Shared Vaults tab: name + storage path (UNC/SMB/local folder).",
        "Saved to shared-vaults.json (DPAPI LocalMachine).",
        "Enterprise Settings → Shared vault combo: pick org vault or default %PROGRAMDATA% path.",
        "Selection persisted in %LOCALAPPDATA%\\FortivaEnterprise\\user.prefs.json.",
        "MaxSeats: seats.dat tracks machineId + userSid; stale seats pruned after 90 days.",
        "New enterprise unlock blocked when active seats >= MaxSeats (existing seat refreshes LastSeen).",
    ])

    add_content_slide(prs, "Audit Logging", [
        "Events: unlock attempt/success/failure, lock, policy violation, snapshot restore, bridge access.",
        "Format: JSONL lines with HMAC-SHA256 signature (AuditIntegrity) — tamper detection.",
        "HMAC key: .audit.hmac.key under audit directory, DPAPI LocalMachine protected.",
        "Personal audit: per-user under FortivaPersonal\\audit.",
        "Enterprise audit: %PROGRAMDATA%\\Fortiva\\audit\\ — exportable JSONL from Admin or client.",
        "Exports contain metadata only — never passwords or master keys.",
    ])

    add_section_slide(prs, "Part 4 — Operations & QA")

    add_code_slide(prs, "Build & Test (Developer / IT Validation)", [
        "# Core + tests (no Visual Studio required):",
        "dotnet build src\\Fortiva.Core\\Fortiva.Core.csproj -c Release",
        "dotnet test tests\\Fortiva.Core.Tests\\Fortiva.Core.Tests.csproj",
        "dotnet test tests\\Fortiva.AppHost.Tests\\Fortiva.AppHost.Tests.csproj",
        "",
        "# Full release + installers:",
        ".\\build-release.ps1 -Version 1.0.x",
        ".\\build-installers.ps1 -Version 1.0.x",
    ])

    add_content_slide(prs, "Import, Export & Backup", [
        "Encrypted export (.fva backup): re-encrypts vault payload with export password (Argon2id + AES-GCM).",
        "Plaintext CSV: allowed on Personal with warning; enterprise policy can block (NoPlaintext).",
        "Import: CSV/KeePass CSV while unlocked; bulk import with row/size limits.",
        "Portable mode (Personal): vault on USB at …\\Fortiva\\vault.fva; startup prompt if drive missing.",
        "Recommendation: periodic encrypted export stored offline separate from master password.",
    ])

    add_content_slide(prs, "Threat Model Summary (What Fortiva Does NOT Claim)", [
        "Does NOT protect against kernel-level malware or memory scraping while vault is unlocked.",
        "Does NOT provide cloud sync, account recovery, or master password reset.",
        "Does NOT send vault data over the network (Personal may HTTPS check for updates only).",
        "DOES protect offline vault at rest with Argon2id + AES-GCM + DPAPI layering.",
        "DOES enforce enterprise policy, licensing, seats, audit, and rollback detection.",
        "Full detail: docs/THREAT-MODEL.md, docs/SECURITY-PENTEST-REPORT.md, docs/FULL-AUDIT-2026.md.",
    ])

    add_title_slide(
        prs,
        "Questions?",
        "Repository: Fortiva  |  Docs: docs/UserManual.md  |  Support: icmclab studio",
    )

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
